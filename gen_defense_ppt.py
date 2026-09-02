# -*- coding: utf-8 -*-
"""生成「智伴成长」答辩 PPT（38 页，16:9，10 分钟版）

按答辩框架要求（创意组）重构：
起源 → 行业/市场 → 实地调研（方法+结果，至少两页）→ 痛点 → 定位 →
产品总述 → 核心技术×3（问题-方案-效果）→ 量化指标 → 研发历程 → 专利布局 →
竞品（差异化/对比/壁垒）→ 商业模式 → 单校账 → 应用案例 → 用户反馈 → 问卷设计 →
财务 → 敏感性 → 融资 → 项目延伸 → 负责人+个人成长 → 成员+团队成长 → 专家顾问 →
学校支撑 → 媒体报道 → 社会价值 → 发展规划 → 风险 → 结语
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn

BLUE = RGBColor(0x1F, 0x4E, 0x79)
MIDBLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHTBLUE = RGBColor(0xDE, 0xEA, 0xF6)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
LIGHTORANGE = RGBColor(0xFB, 0xE5, 0xD6)
GRAY = RGBColor(0x7F, 0x7F, 0x7F)
MGRAY = RGBColor(0xBF, 0xBF, 0xBF)
DARK = RGBColor(0x26, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BGRAY = RGBColor(0xF2, 0xF2, 0xF2)

FONT = '微软雅黑'
TOTAL = 38
OUT = '/Users/xulingexu/Desktop/大创/智伴成长_答辩PPT.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def set_ea(run, name=FONT):
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn('a:latin'))
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        if latin is not None:
            latin.addnext(ea)
        else:
            rPr.append(ea)
    ea.set('typeface', name)


def run_tf(run, text, size=14, bold=False, color=DARK, font=FONT):
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    set_ea(run, font)


def add_text(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (runs, align, sa, ls) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if sa:
            p.space_after = Pt(sa)
        if ls:
            p.line_spacing = ls
        for (t, s, b, c) in runs:
            r = p.add_run()
            run_tf(r, t, s, b, c)
    return tb


def add_shape(slide, st, x, y, w, h, fill=None, line=None, line_w=0.75, radius=None):
    sp = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    if radius is not None and st == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    sp.text_frame.margin_left = sp.text_frame.margin_right = 0
    sp.text_frame.margin_top = sp.text_frame.margin_bottom = 0
    return sp


def card(slide, x, y, w, h, fill=WHITE, line=MGRAY, radius=0.06, line_w=1.0):
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
                     fill=fill, line=line, line_w=line_w, radius=radius)


def shape_text(sp, paras, anchor=MSO_ANCHOR.MIDDLE):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (runs, align, sa, ls) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if sa:
            p.space_after = Pt(sa)
        if ls:
            p.line_spacing = ls
        for (t, s, b, c) in runs:
            r = p.add_run()
            run_tf(r, t, s, b, c)


def new_slide():
    return prs.slides.add_slide(BLANK)


def header(slide, title, idx, accent=ORANGE):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.55, 0.52, 0.13, 0.56, fill=BLUE)
    add_text(slide, 0.85, 0.40, 11.0, 0.8,
             [([(title, 25, True, BLUE)], PP_ALIGN.LEFT, 0, 1.0)])
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.55, 1.18, 2.4, 0.045, fill=accent)
    add_text(slide, 12.35, 7.02, 0.8, 0.35,
             [([(f'{idx:02d} / {TOTAL}', 10, False, GRAY)], PP_ALIGN.RIGHT, 0, 1.0)])


def chip(slide, x, y, w, h, text, fill, tcolor=WHITE, size=13, bold=True):
    sp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
                   fill=fill, radius=0.5)
    shape_text(sp, [([(text, size, bold, tcolor)], PP_ALIGN.CENTER, 0, 1.0)])
    return sp


def tag(slide, x, y, text):
    """右上角"计划/目标/待填"橙色小标签"""
    chip(slide, x, y, 1.15, 0.4, text, ORANGE, size=12)


def psr_cards(slide, x, y, w, h, problem, solution, effect):
    """核心技术页：问题-方案-效果三栏"""
    items = [('原有技术的问题', problem, MGRAY, DARK),
             ('我们的方法', solution, ORANGE, ORANGE),
             ('达到的效果', effect, BLUE, BLUE)]
    for i, (t, d, c, tc) in enumerate(items):
        cx = x + i * (w + 0.2)
        cd = card(slide, cx, y, w, h, fill=WHITE, line=c, line_w=1.5)
        chip(slide, cx + w / 2 - 0.9, y - 0.2, 1.8, 0.4, t, c, size=12)
        add_text(slide, cx + 0.2, y + 0.45, w - 0.4, h - 0.6,
                 [([(ln, 11.5, False, DARK)], PP_ALIGN.CENTER, 0, 1.25) for ln in d.split('\n')])


# ================= P1 封面 =================
s = new_slide()
add_shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.334, 7.5, fill=BLUE)
add_shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.334, 0.12, fill=ORANGE)
add_text(s, 0.9, 0.75, 11.5, 0.5,
         [([('中国国际大学生创新大赛 · 高教主赛道（创意组）', 15, False, RGBColor(0xBF, 0xD3, 0xE6))], PP_ALIGN.LEFT, 0, 1.0)])
add_text(s, 0.9, 2.15, 11.5, 1.5,
         [([('智伴成长', 54, True, WHITE)], PP_ALIGN.LEFT, 0, 1.0)])
add_text(s, 0.95, 3.45, 11.5, 0.6,
         [([('—— 基于 AI 与朋辈导航的大学生全周期智能导学系统', 20, False, RGBColor(0xD9, 0xE2, 0xF3))], PP_ALIGN.LEFT, 0, 1.0)])
add_shape(s, MSO_SHAPE.RECTANGLE, 0.95, 4.25, 2.6, 0.05, fill=ORANGE)
add_text(s, 0.95, 4.5, 11.5, 0.5,
         [([('让每一段大学时光，都有智慧与温暖相伴', 15, False, RGBColor(0xF4, 0xB1, 0x83))], PP_ALIGN.LEFT, 0, 1.0)])
add_text(s, 0.95, 5.9, 11.5, 1.2,
         [([('答辩人：徐麟阁　｜　团队成员：徐麟阁、董芮铭、章宇洲', 13, False, WHITE)], PP_ALIGN.LEFT, 0, 1.2),
          ([('指导教师：陈圣波　｜　上海大学计算机工程与科学学院', 13, False, WHITE)], PP_ALIGN.LEFT, 0, 1.2),
          ([('2026 年 9 月', 13, False, RGBColor(0xBF, 0xD3, 0xE6))], PP_ALIGN.LEFT, 0, 1.2)])

# ================= P2 目录 =================
s = new_slide()
header(s, '目录', 2)
toc = [
    ('01', '起源与调研', '项目起源 · 行业市场 · 实地调研 · 痛点'),
    ('02', '产品与核心技术', '定位 · 产品 · 三大核心技术 · 指标'),
    ('03', '竞品与商业模式', '差异化 · 竞品对比 · 商业模式'),
    ('04', '应用与财务', '案例 · 反馈 · 财务 · 融资'),
    ('05', '团队与成长', '负责人 · 成员 · 顾问 · 学校支撑'),
    ('06', '传播与社会价值', '媒体报道 · 社会价值'),
    ('07', '规划与展望', '发展规划 · 风险 · 结语'),
]
positions = [(0.9, 1.9), (6.8, 1.9), (0.9, 3.3), (6.8, 3.3), (0.9, 4.7), (6.8, 4.7), (0.9, 6.1)]
for (num, t, sub), (x, y) in zip(toc, positions):
    card(s, x, y, 5.5, 1.05, fill=WHITE, line=MGRAY)
    ov = add_shape(s, MSO_SHAPE.OVAL, x + 0.25, y + 0.27, 0.5, 0.5, fill=BLUE)
    shape_text(ov, [([(num, 15, True, WHITE)], PP_ALIGN.CENTER, 0, 1.0)])
    add_text(s, x + 0.95, y + 0.16, 4.4, 0.8,
             [([(t, 16, True, DARK)], PP_ALIGN.LEFT, 0, 1.1),
              ([(sub, 11.5, False, GRAY)], PP_ALIGN.LEFT, 0, 1.1)])

# ================= P3 项目起源 =================
s = new_slide()
header(s, '项目起源：政策 · 学校 · 个人三重驱动', 3)
srcs = [
    ('政策驱动', '教育数字化战略行动与"人工智能+"行动明确推动智能技术融入校园场景', BLUE),
    ('学校培养', '机器学习、自然语言处理课程 + 三项实训（Word2Vec / TextCNN / 知识图谱）奠定技术底座', MIDBLUE),
    ('个人经历', '团队成员亲历新生信息焦虑与朋辈经验"口口相传"的断层，萌生"数字学长"构想', ORANGE),
]
for i, (t, d, c) in enumerate(srcs):
    x = 0.9 + i * 3.95
    cd = card(s, x, 1.7, 3.6, 3.6, fill=WHITE, line=c, line_w=1.5)
    chip(s, x + 1.28, 2.0, 1.04, 0.44, t[:2], c, size=13)
    add_text(s, x + 0.25, 2.7, 3.1, 2.4,
             [([(t, 16, True, DARK)], PP_ALIGN.CENTER, 4, 1.1),
              ([(d, 12, False, GRAY)], PP_ALIGN.CENTER, 0, 1.3)])
add_text(s, 0.9, 5.7, 11.5, 0.8,
         [([('由此立项：用 AI 把学长经验变成"数字学长"，服务每一位大学生的全周期成长', 13.5, True, BLUE)], PP_ALIGN.CENTER, 0, 1.1)])

# ================= P4 行业背景 =================
s = new_slide()
header(s, '行业背景：校园智能服务风口（PEST）', 4)
pest = [
    ('P 政策', '教育数字化、AI+ 行动，高校智能服务是明确方向', BLUE, LIGHTBLUE),
    ('E 经济', '高校信息化预算逐年增长，年费订阅模式可核算', MIDBLUE, LIGHTBLUE),
    ('S 社会', '00 后数字原生习惯；心理健康受关注；朋辈互助数字化空白', ORANGE, LIGHTORANGE),
    ('T 技术', '大模型 / RAG / 知识图谱成熟开源，低成本落地可行', BLUE, LIGHTBLUE),
]
for i, (t, d, c, fill) in enumerate(pest):
    x = 0.9 + (i % 2) * 5.95
    y = 1.7 + (i // 2) * 2.3
    cd = card(s, x, y, 5.6, 1.9, fill=fill, line=None)
    add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 0.12, 1.9, fill=c)
    add_text(s, x + 0.4, y + 0.2, 5.0, 1.5,
             [([(t, 16, True, c)], PP_ALIGN.LEFT, 4, 1.1),
              ([(d, 12, False, DARK)], PP_ALIGN.LEFT, 0, 1.25)])
add_text(s, 0.9, 6.45, 11.5, 0.5,
         [([('结论：校园 AI 助手赛道已被头部厂商验证，但"成长陪伴"场景仍处空白', 13, True, BLUE)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P5 市场规模 =================
s = new_slide()
header(s, '市场规模：TAM / SAM / SOM 三层漏斗', 5)
funnel = [
    ('TAM　总体市场', '全国 3,000+ 所高校 × 年均 5 万元', '约 15,000 万元/年', 11.4, MIDBLUE, WHITE),
    ('SAM　可服务市场', '具备采购意愿的本科院校约 800 所（约 27%）', '约 4,000 万元/年', 9.4, BLUE, WHITE),
    ('SOM　可获取市场', '第 3 年累计签约 40 所 × 加权年费约 5 万元', '约 200 万元/年', 7.4, ORANGE, WHITE),
]
y = 1.75
for i, (name, basis, val, w, c, tc) in enumerate(funnel):
    x = (13.333 - w) / 2
    bar = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 1.15, fill=c, radius=0.12)
    shape_text(bar, [([(name, 15, True, tc), ('　' + val, 15, True, WHITE)], PP_ALIGN.CENTER, 0, 1.1),
                     ([(basis, 11.5, False, RGBColor(0xEA, 0xEF, 0xF6))], PP_ALIGN.CENTER, 0, 1.1)])
    if i < 2:
        add_shape(s, MSO_SHAPE.DOWN_ARROW, (13.333 - 0.32) / 2, y + 1.13, 0.32, 0.32, fill=GRAY)
    y += 1.5
note = card(s, 9.8, 6.35, 2.35, 0.8, fill=LIGHTBLUE, line=None)
shape_text(note, [([('约占 SAM 5%', 13, True, BLUE), ('即可达成目标', 11, False, DARK)], PP_ALIGN.CENTER, 0, 1.15)])
add_text(s, 0.9, 6.45, 8.6, 0.7,
         [([('在校生约 3,800 万人 · 每年新增新生 1,000 万+ · 朋辈导学场景供给空白', 12.5, False, DARK)], PP_ALIGN.LEFT, 0, 1.1)])

# ================= P6 实地调研·方法 =================
s = new_slide()
header(s, '实地调研（1/2）：调研方法', 6)
methods = [
    ('问卷调查', '大一至大四全覆盖，目标样本 ≥300 份，线上发放、全程留存原始答卷', BLUE),
    ('半结构化访谈', '学工教师 5—10 人：事务性答疑负担、减负意愿、采购决策链', MIDBLUE),
    ('桌面研究', '教育部政策文件；竞品上市公司年报与官网（新开普、金智教育）', ORANGE),
]
for i, (t, d, c) in enumerate(methods):
    x = 0.9 + i * 3.95
    cd = card(s, x, 1.7, 3.6, 2.9, fill=WHITE, line=c, line_w=1.5)
    chip(s, x + 1.28, 2.0, 1.04, 0.44, t[:2], c, size=13)
    add_text(s, x + 0.25, 2.7, 3.1, 1.8,
             [([(t, 15, True, DARK)], PP_ALIGN.CENTER, 4, 1.1),
              ([(d, 11.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.3)])
ban = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 4.95, 11.5, 1.15, fill=LIGHTBLUE, radius=0.12)
shape_text(ban, [([('执行安排：2026 年 9 月启动 → 10 月底输出调研报告 → 校准产品功能优先级与财务假设', 13.5, True, BLUE)], PP_ALIGN.CENTER, 0, 1.1),
                 ([('调研全程留痕：原始问卷、访谈记录、现场台账、数据清洗脚本，满足 2026 评审真实性要求', 11.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.1)])
add_text(s, 0.9, 6.4, 11.5, 0.5,
         [([('问卷完整设计见商业计划书附录一（8 道核心题）', 11.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P7 实地调研·过程与结果 =================
s = new_slide()
header(s, '实地调研（2/2）：过程与结果', 7)
c1 = card(s, 0.9, 1.6, 5.6, 4.2, fill=LIGHTBLUE, line=None)
add_text(s, 1.25, 1.85, 5.0, 0.5, [([('已完成：三项验证基础', 15, True, BLUE)], PP_ALIGN.LEFT, 0, 1.0)])
add_text(s, 1.25, 2.5, 5.0, 3.1,
         [([('① 团队切身体验：', 12, True, DARK), ('全体成员亲历新生适应、选课、备考全周期，对痛点有第一手观察', 12, False, DARK)], PP_ALIGN.LEFT, 6, 1.25),
          ([('② 政策与配比数据：', 12, True, DARK), ('教育部规定辅导员按不低于 1:200 师生比配置，个性化需求存在结构性缺口', 12, False, DARK)], PP_ALIGN.LEFT, 6, 1.25),
          ([('③ 行业动向佐证：', 12, True, DARK), ('金智教育、新开普 2025—2026 密集上线校园 AI 助手，需求被供给端验证', 12, False, DARK)], PP_ALIGN.LEFT, 0, 1.25)])
c2 = card(s, 6.85, 1.6, 5.6, 4.2, fill=LIGHTORANGE, line=None)
add_text(s, 7.2, 1.85, 5.0, 0.5, [([('进行中：正式调研（2026.09—10）', 15, True, ORANGE)], PP_ALIGN.LEFT, 0, 1.0)])
add_text(s, 7.2, 2.5, 5.0, 3.1,
         [([('① 9 月上旬：', 12, True, ORANGE), ('问卷投放（≥300 份）+ 学工教师访谈启动', 12, False, DARK)], PP_ALIGN.LEFT, 6, 1.25),
          ([('② 9 月下旬：', 12, True, ORANGE), ('访谈 5—10 人完成，整理记录与台账', 12, False, DARK)], PP_ALIGN.LEFT, 6, 1.25),
          ([('③ 10 月：', 12, True, ORANGE), ('数据清洗与假设检验（H1—H4）', 12, False, DARK)], PP_ALIGN.LEFT, 6, 1.25),
          ([('④ 10 月底：', 12, True, ORANGE), ('调研报告 → 校准功能优先级与财务假设', 12, False, DARK)], PP_ALIGN.LEFT, 0, 1.25)])
add_text(s, 0.9, 6.15, 11.5, 0.7,
         [([('诚实声明：正式调研结果未出前，本答辩所有结论均标注为"已验证基础 + 待验证假设"，不引用未执行的数据', 12, True, BLUE)], PP_ALIGN.CENTER, 0, 1.1)])

# ================= P8 痛点分析 =================
s = new_slide()
header(s, '痛点分析：学生端与学校端的双向困境', 8)
add_text(s, 0.9, 1.45, 5.6, 0.5, [([('学生端之痛', 17, True, BLUE)], PP_ALIGN.CENTER, 0, 1.0)])
add_text(s, 6.85, 1.45, 5.6, 0.5, [([('学校端之痛', 17, True, ORANGE)], PP_ALIGN.CENTER, 0, 1.0)])
stu = [('信息获取碎片化', '办事指南散落在官网、部门与聊天记录中'),
       ('选课决策迷茫', '学长经验靠口口相传，缺乏系统化渠道'),
       ('心理疏导滞后', '愿意倾诉、不愿当面求助')]
sch = [('辅导员“一对多”', '师生比不低于 1:200，难以兼顾个体需求'),
       ('经验难以沉淀', '学长经验随毕业流失，无法代际传承'),
       ('重复答疑负担重', '事务性咨询占用大量学工工时')]
c1 = card(s, 0.9, 2.0, 5.6, 3.9, fill=LIGHTBLUE, line=None)
c2 = card(s, 6.85, 2.0, 5.6, 3.9, fill=LIGHTORANGE, line=None)
for (cd, items, c) in [(c1, stu, BLUE), (c2, sch, ORANGE)]:
    paras = []
    for t, d in items:
        paras.append(([( '▍', 14, True, c), (t, 14.5, True, DARK)], PP_ALIGN.LEFT, 2, 1.15))
        paras.append(([(d, 12, False, GRAY)], PP_ALIGN.LEFT, 14, 1.2))
    shape_text(cd, paras, anchor=MSO_ANCHOR.MIDDLE)
    cd.text_frame.margin_left = Inches(0.35)
    cd.text_frame.margin_right = Inches(0.25)
ban = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 6.05, 11.5, 1.0, fill=WHITE, line=ORANGE, line_w=1.5, radius=0.12)
shape_text(ban, [([('双向痛点 → 一个方案：', 13.5, True, DARK), ('“智伴成长”数字学长', 13.5, True, ORANGE)], PP_ALIGN.CENTER, 1, 1.1),
                 ([('调研留痕：2026.09—10 问卷 ≥300 份 + 学工访谈，原始问卷/访谈记录/现场台账全程留存（可核验）', 10.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.1)])

# ================= P9 核心定位 =================
s = new_slide()
header(s, '核心定位：什么是“数字学长”？', 9)
ban = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 1.6, 11.5, 1.15, fill=LIGHTBLUE, radius=0.12)
shape_text(ban, [([('一句话定义：', 17, True, BLUE), ('为每一位大学生配备的 7×24 在线智能导学助手——', 17, True, DARK)], PP_ALIGN.CENTER, 0, 1.15),
                 ([('将 AI 的高效与朋辈导师的亲和力结合，覆盖大学全周期成长', 14, False, GRAY)], PP_ALIGN.CENTER, 0, 1.15)])
roles = [
    ('工具', '高效的知识管理工具', '“查什么、怎么办、去哪办”，三步之内得到答案', MIDBLUE),
    ('导师', '可靠的学业与生活导师', '选课建议 · 挂科预警 · 关键节点主动提醒', BLUE),
    ('陪伴', '有温度的陪伴者', '情绪感知与心理支持 · 危机信号自动转介', ORANGE),
]
for i, (tg, t, d, c) in enumerate(roles):
    x = 0.9 + i * 3.95
    card(s, x, 3.25, 3.6, 2.6, fill=WHITE, line=c, line_w=1.5)
    chip(s, x + 1.28, 3.5, 1.04, 0.44, tg, c, size=13)
    add_text(s, x + 0.3, 4.15, 3.0, 1.6,
             [([(t, 16, True, DARK)], PP_ALIGN.CENTER, 4, 1.15),
              ([(d, 12, False, GRAY)], PP_ALIGN.CENTER, 0, 1.25)])
add_text(s, 0.9, 6.35, 11.5, 0.5,
         [([('三重角色 = 工具 + 导师 + 陪伴，覆盖“适应—成长—关怀”完整用户旅程', 13, True, BLUE)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P10 产品总述 =================
s = new_slide()
header(s, '产品总述：四层技术架构', 10)
layers = [
    ('01', '数据层', '官网公告 · 办事流程 · 课程资料 · 学长经验 → 清洗去重', LIGHTBLUE, BLUE),
    ('02', '知识层', 'Neo4j 知识图谱 + 向量数据库——“数字学长的大脑”', RGBColor(0xBD, 0xD7, 0xEE), BLUE),
    ('03', '服务层', 'RAG 问答引擎 · 个性化推荐 · 上下文感知（情绪识别）', LIGHTBLUE, BLUE),
    ('04', '交互层', 'Electron 跨平台客户端 · 文字/语音交互 · 主动推送', LIGHTBLUE, BLUE),
]
y = 1.6
for num, name, desc, fill, c in layers:
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.5, y, 10.3, 0.92, fill=fill, radius=0.14)
    sq = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.75, y + 0.21, 0.52, 0.52, fill=c, radius=0.2)
    shape_text(sq, [([(num, 14, True, WHITE)], PP_ALIGN.CENTER, 0, 1.0)])
    add_text(s, 2.55, y + 0.13, 9.0, 0.7,
             [([(name, 16, True, BLUE), ('　' + desc, 12.5, False, DARK)], PP_ALIGN.LEFT, 0, 1.05)])
    if num != '04':
        add_shape(s, MSO_SHAPE.DOWN_ARROW, 6.52, y + 0.9, 0.3, 0.26, fill=GRAY)
    y += 1.2
add_text(s, 1.5, 6.55, 10.3, 0.5,
         [([('技术底座：RAG 约束生成（答案可溯源） + 知识图谱关联推荐 + TextCNN 情绪识别', 12.5, True, BLUE)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P11 四大功能模块 =================
s = new_slide()
header(s, '四大功能模块：覆盖大学全周期', 11)
mods = [
    ('生活导航', '校园事务实操指引', '校园卡激活 · 宿舍水电 · 图书馆预约', MIDBLUE),
    ('学业规划', '选课与成长指导', '选课组合建议 · 挂科风险提醒 · 学习转型指导', BLUE),
    ('心理疏导', '情绪感知与支持', '情绪识别 · 心理调节方案 · 危机转介提醒', ORANGE),
    ('校园融入', '归属感建设', '社团推荐 · 校园文化传递 · 归属感建设', MIDBLUE),
]
for i, (t, sub, d, c) in enumerate(mods):
    x = 0.9 + (i % 2) * 5.95
    y = 1.65 + (i // 2) * 2.15
    card(s, x, y, 5.6, 1.85, fill=WHITE, line=c, line_w=1.5)
    add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 0.12, 1.85, fill=c)
    add_text(s, x + 0.4, y + 0.22, 5.0, 1.4,
             [([(t, 17, True, c), ('　' + sub, 13, False, DARK)], PP_ALIGN.LEFT, 4, 1.1),
              ([(d, 12, False, GRAY)], PP_ALIGN.LEFT, 0, 1.2)])
ban = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 6.15, 11.5, 0.75, fill=LIGHTORANGE, radius=0.5)
shape_text(ban, [([('+ 个性化推荐与主动提醒：', 14, True, ORANGE),
                   ('选课截止 · 考试周 · 新生报到等 20+ 类关键节点主动推送', 14, True, DARK)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P12-14 核心技术 =================
s = new_slide()
header(s, '核心技术（一）：RAG 智能问答', 12)
psr_cards(s, 0.9, 2.4, 3.7, 3.4,
          '通用大模型缺乏校园私域知识\n答案不可溯源\n校园场景易产生"幻觉"',
          'RAG 检索增强生成：\n理解问题 → 检索知识库\n→ 约束生成回答\n全部答案附来源引用',
          '答案 100% 可溯源\n高频问答准确率\n≥95%（目标值）\n白名单答案库兜底')
add_text(s, 0.9, 6.3, 11.5, 0.6,
         [([('技术来源：团队 Word2Vec 实训的语义匹配能力迁移至检索排序环节', 12, False, GRAY)], PP_ALIGN.CENTER, 0, 1.0)])

s = new_slide()
header(s, '核心技术（二）：Neo4j 校园知识图谱', 13)
psr_cards(s, 0.9, 2.4, 3.7, 3.4,
          '选课/社团/流程等关联信息散落\n传统关系型存储\n无法表达实体间复杂关系',
          'Neo4j 图数据库构建\n人、事、地、物实体图谱\n多跳查询与关联推荐\n学长经验众包沉淀',
          '选课组合建议\n挂科风险提醒\n社团匹配推荐\n知识条目 ≥2,000（目标）')
add_text(s, 0.9, 6.3, 11.5, 0.6,
         [([('技术来源：团队医学知识图谱实训的实体抽取与关系构建方法论直接复用', 12, False, GRAY)], PP_ALIGN.CENTER, 0, 1.0)])

s = new_slide()
header(s, '核心技术（三）：主动提醒与情绪感知', 14)
psr_cards(s, 0.9, 2.4, 3.7, 3.4,
          '传统问答被动响应\n缺乏对关键时间节点\n与用户情绪的感知能力',
          '上下文感知引擎\nTextCNN 情绪三分类\n20+ 类场景提醒引擎\n危机信号转介机制',
          '负面情绪 5 分钟内\n推送心理调节方案\n选课截止/考试周\n主动提醒（目标值）')
add_text(s, 0.9, 6.3, 11.5, 0.6,
         [([('技术来源：TextCNN 文本分类实训直接转化为情绪识别模块', 12, False, GRAY)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P15 核心量化指标 =================
s = new_slide()
header(s, '核心量化指标：差异化必须可验证', 15)
stats = [
    ('≥95%', '高频问答准确率', 'Top 200 问题双盲抽检', BLUE),
    ('100%', '答案附来源可溯源', '系统强制机制', BLUE),
    ('20+ 类', '主动提醒场景', '选课/考试周/报到等', ORANGE),
    ('≥85%', '情绪识别准确率', '负面情绪 5 分钟内关怀', ORANGE),
    ('≥40%', '试点校月活跃率', '次月留存 ≥60%', BLUE),
    ('≤2 秒', '文字问答首字响应', '客户端埋点统计', BLUE),
]
for i, (num, t, d, c) in enumerate(stats):
    x = 0.9 + (i % 3) * 3.95
    y = 1.7 + (i // 3) * 2.15
    card(s, x, y, 3.6, 1.85, fill=WHITE, line=MGRAY)
    add_text(s, x + 0.25, y + 0.22, 3.1, 1.4,
             [([(num, 30, True, c)], PP_ALIGN.LEFT, 2, 1.0),
              ([(t, 13.5, True, DARK)], PP_ALIGN.LEFT, 2, 1.1),
              ([(d, 11, False, GRAY)], PP_ALIGN.LEFT, 0, 1.1)])
add_text(s, 0.9, 6.3, 11.5, 0.6,
         [([('口径说明：以上为 MVP/V1.0 验收目标值 + 验证方式，试点实测校准；', 12, False, GRAY),
            ('未达标不对外宣传', 12, True, ORANGE)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P16 研发历程 =================
s = new_slide()
header(s, '研发历程：从实训到产品的迭代之路', 16)
steps = [
    ('2026.01—07', '实训积累', 'Word2Vec / TextCNN / 医学知识图谱三项实训', BLUE),
    ('2026.08', '开题立项', '开题报告、计划书、商业计划书、答辩 PPT', MIDBLUE),
    ('2026.09—10', '调研与架构', '问卷访谈调研；四层架构设计；技术选型', ORANGE),
    ('2026.11—2027.01', 'MVP 开发', 'RAG 问答链路；基础知识库；生活导航模块', BLUE),
    ('2027.02—04', 'V1.0 迭代', '知识图谱上线；学业规划；跨平台客户端', MIDBLUE),
    ('2027.05—06', '试点与结题', '校内全量部署；成效报告；结题验收', ORANGE),
]
for i, (t, name, d, c) in enumerate(steps):
    x = 0.9 + (i % 3) * 3.95
    y = 1.6 + (i // 3) * 2.25
    card(s, x, y, 3.6, 1.9, fill=WHITE, line=c, line_w=1.5)
    add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 0.12, 1.9, fill=c)
    add_text(s, x + 0.35, y + 0.15, 3.1, 1.6,
             [([(t, 12, True, c)], PP_ALIGN.LEFT, 3, 1.1),
              ([(name, 15, True, DARK)], PP_ALIGN.LEFT, 3, 1.1),
              ([(d, 11, False, GRAY)], PP_ALIGN.LEFT, 0, 1.15)])
add_text(s, 0.9, 6.35, 11.5, 0.5,
         [([('每一步均留存过程材料：代码提交记录、实验笔记、文档版本，可追溯可核验', 12, False, GRAY)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P17 知识产权布局 =================
s = new_slide()
header(s, '知识产权布局（计划）', 17)
tag(s, 11.9, 1.15, '计划')
ips = [
    ('软件著作权', '目标 1 项：2027 年 1 月前提交"智伴成长智能导学系统"软著申请', BLUE),
    ('查新报告', '2026 年 12 月完成技术查新，明确与同类产品的差异化技术特征', MIDBLUE),
    ('代码与文档沉淀', 'GitHub 仓库版本化留存（已执行）；实验笔记与文档全程归档', ORANGE),
]
for i, (t, d, c) in enumerate(ips):
    x = 0.9 + i * 3.95
    cd = card(s, x, 1.7, 3.6, 2.9, fill=WHITE, line=c, line_w=1.5)
    chip(s, x + 1.28, 2.0, 1.04, 0.44, t[:2], c, size=13)
    add_text(s, x + 0.25, 2.7, 3.1, 1.8,
             [([(t, 15, True, DARK)], PP_ALIGN.CENTER, 4, 1.1),
              ([(d, 11.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.3)])
add_text(s, 0.9, 5.0, 11.5, 1.6,
         [([('成果边界声明：', 13, True, BLUE), ('项目核心技术（RAG 链路、知识图谱、跨平台客户端）均由团队学生独立完成，', 12.5, False, DARK)], PP_ALIGN.CENTER, 2, 1.25),
          ([('指导教师负责技术把关与资源协调；不存在将教师成果包装为学生成果的情形', 12.5, False, DARK)], PP_ALIGN.CENTER, 0, 1.25),
          ([('（软著与查新均为执行计划，将在时间表节点落实）', 11, True, ORANGE)], PP_ALIGN.CENTER, 0, 1.2)])

# ================= P18 差异化 =================
s = new_slide()
header(s, '差异化：错位竞争——“办事效率 vs 成长陪伴”', 18)
add_text(s, 0.9, 1.45, 5.7, 0.5, [([('竞品主场：办事效率', 16, True, GRAY)], PP_ALIGN.CENTER, 0, 1.0)])
add_text(s, 6.75, 1.45, 5.7, 0.5, [([('智伴成长：成长陪伴', 16, True, ORANGE)], PP_ALIGN.CENTER, 0, 1.0)])
c1 = card(s, 0.9, 2.0, 5.7, 3.75, fill=BGRAY, line=None)
shape_text(c1, [([('一卡通服务 · 事务办理 · 官方资讯问答', 14, True, DARK)], PP_ALIGN.CENTER, 6, 1.2),
                ([('金智教育（1,200+ 所高校）', 12.5, False, GRAY)], PP_ALIGN.CENTER, 2, 1.2),
                ([('新开普（一卡通约 45% 份额）', 12.5, False, GRAY)], PP_ALIGN.CENTER, 2, 1.2),
                ([('高校自研 AI 助手（清华/浙大等）', 12.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.2)], anchor=MSO_ANCHOR.MIDDLE)
c2 = card(s, 6.75, 2.0, 5.7, 3.75, fill=LIGHTORANGE, line=ORANGE, line_w=1.5)
shape_text(c2, [([('选课避雷 · 学习方法 · 心理经验 · 社团融入', 14, True, ORANGE)], PP_ALIGN.CENTER, 6, 1.2),
                ([('朋辈经验众包沉淀，代际传承', 12.5, False, DARK)], PP_ALIGN.CENTER, 2, 1.2),
                ([('“数字学长”人格 + 情绪感知', 12.5, False, DARK)], PP_ALIGN.CENTER, 2, 1.2),
                ([('无人系统化覆盖的空白地带', 12.5, False, DARK)], PP_ALIGN.CENTER, 0, 1.2)], anchor=MSO_ANCHOR.MIDDLE)
ov = add_shape(s, MSO_SHAPE.OVAL, 5.79, 3.3, 1.75, 1.05, fill=WHITE, line=ORANGE, line_w=1.5)
shape_text(ov, [([('错位', 16, True, ORANGE)], PP_ALIGN.CENTER, 0, 1.0)])
add_text(s, 0.9, 6.15, 11.5, 0.7,
         [([('策略：合作大于对抗（直销 + 被集成） · 单点极致 · 窗口期 1—2 年', 13.5, True, BLUE)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P19 竞品对比 =================
s = new_slide()
header(s, '竞品对比：错位优势一目了然', 19)
rows, cols = 6, 5
tbl_shape = s.shapes.add_table(rows, cols, Inches(0.9), Inches(1.55), Inches(11.5), Inches(5.0))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(1.7)
for j in range(1, 5):
    tbl.columns[j].width = Inches(2.45)
tbl.rows[0].height = Inches(0.6)
for r in range(1, rows):
    tbl.rows[r].height = Inches(0.8)
data = [
    ['维度', '智伴成长', '厂商 AI 助手', '高校自研', '通用大模型'],
    ['产品形态', '数字学长 · 全周期导学', '服务 App + 办事助手', '校内问答机器人', '通用对话助手'],
    ['朋辈经验', '众包沉淀 · 代际传承', '基本不涉及', '基本不涉及', '无'],
    ['主动提醒', '20+ 类关键场景', '服务办理类为主', '较少', '无'],
    ['情感角色', '学长人设 + 情绪感知', '数字人客服（工具）', '工具属性', '通用人格'],
    ['收费模式', 'B2B 年订阅', '项目制 / 平台费', '校方投入', '免费 / 订阅'],
]
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
            run_tf(run, data[r][c], 13, True, WHITE)
        elif c == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHTBLUE
            run_tf(run, data[r][c], 12, True, DARK)
        elif c == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHTORANGE
            run_tf(run, data[r][c], 12, True, DARK)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            run_tf(run, data[r][c], 11.5, False, GRAY)
add_text(s, 0.9, 6.85, 11.5, 0.4,
         [([('注：竞品信息依据公开资料整理（截至 2026 年 8 月）；量化目标对比见 P15。', 10.5, False, GRAY)], PP_ALIGN.LEFT, 0, 1.0)])

# ================= P20 竞争壁垒 =================
s = new_slide()
header(s, '竞争壁垒：三重壁垒护城河', 20)
walls = [
    ('01', '数据壁垒', '校园私域知识库 + 朋辈经验众包生态\n——厂商与自研系统均不覆盖的数据形态', BLUE, LIGHTBLUE),
    ('02', '场景壁垒', '嵌入选课、报到、考试周关键节点\n“用得上、离不开”的使用习惯', ORANGE, LIGHTORANGE),
    ('03', '角色壁垒', '“数字学长”人格与情感联结\n工具型竞品难以模仿的软性壁垒', MIDBLUE, LIGHTBLUE),
]
for i, (num, t, d, c, fill) in enumerate(walls):
    x = 0.9 + i * 3.95
    card(s, x, 1.7, 3.6, 3.9, fill=fill, line=None)
    add_shape(s, MSO_SHAPE.OVAL, x + 1.43, 2.0, 0.75, 0.75, fill=c)
    add_text(s, x + 1.43, 2.0, 0.75, 0.75, [([(num, 20, True, WHITE)], PP_ALIGN.CENTER, 0, 1.0)])
    add_text(s, x + 0.25, 2.95, 3.1, 0.6, [([(t, 18, True, c)], PP_ALIGN.CENTER, 0, 1.0)])
    add_text(s, x + 0.25, 3.65, 3.1, 1.6, [([(ln, 12, False, DARK)], PP_ALIGN.CENTER, 0, 1.25) for ln in d.split('\n')])
add_text(s, 0.9, 6.05, 11.5, 0.7,
         [([('打法：错位聚焦 · 单点极致 · 以快打慢 —— 1—2 年窗口期内完成 3 校标杆与 40 所目标', 13.5, True, BLUE)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P21 商业模式 =================
s = new_slide()
header(s, '商业模式：B2B 订阅制，学校付费、学生免费', 21)
ban = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 1.5, 11.5, 0.7, fill=BLUE, radius=0.5)
shape_text(ban, [([('B2B2C：', 14, True, WHITE), ('高校按年订阅（学生端完全免费） · 辅导员减负是学校的刚性需求', 13.5, False, WHITE)], PP_ALIGN.CENTER, 0, 1.0)])
prices = [
    ('基础版', '2.98 万元/年', '万人以下院校 / 二级学院\nRAG 问答 + 校园知识库', False),
    ('标准版', '4.98 万元/年', '普通本科院校\n+ 知识图谱 · 学业规划 · 主动提醒', True),
    ('旗舰版', '8.8 万元/年', '重点高校 / 多校区\n+ 个性化推荐 · 私有化部署', False),
]
for i, (t, p, d, hot) in enumerate(prices):
    x = 0.9 + i * 3.95
    card(s, x, 2.55, 3.6, 2.7, fill=WHITE, line=(ORANGE if hot else MGRAY), line_w=(2.0 if hot else 1.0))
    if hot:
        chip(s, x + 2.45, 2.38, 1.0, 0.4, '主打', ORANGE, size=12)
    add_text(s, x + 0.3, 2.85, 3.0, 0.5, [([(t, 17, True, BLUE)], PP_ALIGN.CENTER, 0, 1.0)])
    add_text(s, x + 0.3, 3.45, 3.0, 0.6, [([(p, 22, True, ORANGE)], PP_ALIGN.CENTER, 0, 1.0)])
    add_text(s, x + 0.3, 4.15, 3.0, 1.0, [([(ln, 11.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.2) for ln in d.split('\n')])
add_text(s, 0.9, 5.65, 11.5, 1.0,
         [([('销售结构：基础版 30% / 标准版 55% / 旗舰版 15% → 加权平均年费 4.95 万元', 13, True, DARK)], PP_ALIGN.CENTER, 2, 1.15),
          ([('首年另收实施与数据初始化费 0.5—2 万元（校均 0.8 万元） · 续费率假设 80%', 12, False, GRAY)], PP_ALIGN.CENTER, 0, 1.15)])

# ================= P22 单校经济账 =================
s = new_slide()
header(s, '单校经济账：首年即回本', 22)
add_text(s, 1.0, 1.45, 5.3, 0.5, [([('首年（含实施费）', 16, True, BLUE)], PP_ALIGN.CENTER, 0, 1.0)])
add_text(s, 7.0, 1.45, 5.3, 0.5, [([('续费年（第 2 年起）', 16, True, ORANGE)], PP_ALIGN.CENTER, 0, 1.0)])
rows = [('合同收入', 5.75, 4.95), ('单校毛利', 3.05, 3.05), ('全成本净贡献', 0.75, 1.55)]
UNIT = 0.62
y = 2.15
for name, v1, v2 in rows:
    add_text(s, 1.0, y, 1.7, 0.4, [([(name, 12.5, True, DARK)], PP_ALIGN.RIGHT, 0, 1.0)])
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 2.85, y, max(0.25, v1 * UNIT), 0.42, fill=BLUE, radius=0.5)
    add_text(s, 2.95 + v1 * UNIT, y - 0.02, 1.6, 0.45, [([(f'{v1:g} 万', 12, True, BLUE)], PP_ALIGN.LEFT, 0, 1.0)])
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 7.85, y, max(0.25, v2 * UNIT), 0.42, fill=ORANGE, radius=0.5)
    add_text(s, 7.95 + v2 * UNIT, y - 0.02, 1.6, 0.45, [([(f'{v2:g} 万', 12, True, ORANGE)], PP_ALIGN.LEFT, 0, 1.0)])
    y += 0.78
add_text(s, 1.0, y + 0.05, 11.3, 0.4,
         [([('毛利率：首年 53% → 续费年边际毛利率 62%（客户存续越久，利润越厚）', 12.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.0)])
ban = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 6.15, 11.5, 0.85, fill=ORANGE, radius=0.5)
shape_text(ban, [([('核心结论：', 15, True, WHITE),
                   ('单校首年即回本 · LTV/CAC ≈ 3.7（健康线为 3）· 公司盈亏平衡点 15—16 所', 15, True, WHITE)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P23 应用案例·试点规划 =================
s = new_slide()
header(s, '应用案例：三层试点规划（目标）', 23)
tag(s, 11.9, 1.15, '目标')
cases = [
    ('案例一', '本校全量部署', '校级 · 学生工作部门\n新生入学教育嵌入\n2027.03 目标', BLUE),
    ('案例二', '学院付费试点', '院级 · 二级学院\n标准版订阅\n2027.05 目标', ORANGE),
    ('案例三', '省内 3 校复制', '校级 · 信息化部门\n3 校案例包输出\n2027.06 起目标', MIDBLUE),
]
for i, (t, sub, d, c) in enumerate(cases):
    x = 0.9 + i * 3.95
    card(s, x, 1.8, 3.6, 3.3, fill=WHITE, line=c, line_w=1.5)
    chip(s, x + 1.28, 2.1, 1.04, 0.44, t[:2], c, size=13)
    add_text(s, x + 0.25, 2.75, 3.1, 2.2,
             [([(sub, 15, True, DARK)], PP_ALIGN.CENTER, 4, 1.1)] +
             [([(ln, 11.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.3) for ln in d.split('\n')])
add_text(s, 0.9, 5.5, 11.5, 1.2,
         [([('说明：应用案例为按里程碑规划的试点目标，尚未发生；', 12, False, GRAY),
            ('MVP 完成后按"本校 → 学院 → 省内"逐层推进并留存真实案例材料', 12, True, BLUE)], PP_ALIGN.CENTER, 0, 1.3)])

# ================= P24 用户反馈·内测计划 =================
s = new_slide()
header(s, '用户反馈：内测计划与反馈机制（计划）', 24)
tag(s, 11.9, 1.15, '计划')
fbs = [
    ('500 人内测', '2027.01 校内招募 500 人规模内测，覆盖大一至大四', BLUE),
    ('周度复盘', '点踩率、未解决问题 Top10、满意度变化逐周跟踪', ORANGE),
    ('效果目标', '满意度 ≥4.5/5；月活跃率 ≥40%；次月留存 ≥60%', MIDBLUE),
    ('反馈闭环', '用户反馈沉淀为知识库更新与迭代排期的直接依据', BLUE),
]
for i, (t, d, c) in enumerate(fbs):
    x = 0.9 + (i % 2) * 5.95
    y = 1.7 + (i // 2) * 2.3
    card(s, x, y, 5.6, 1.9, fill=WHITE, line=c, line_w=1.5)
    add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 0.12, 1.9, fill=c)
    add_text(s, x + 0.4, y + 0.2, 5.0, 1.5,
             [([(t, 16, True, c)], PP_ALIGN.LEFT, 4, 1.1),
              ([(d, 12, False, DARK)], PP_ALIGN.LEFT, 0, 1.25)])
add_text(s, 0.9, 6.45, 11.5, 0.5,
         [([('说明：内测与反馈机制为执行计划，正式反馈数据将在 2027.01 起沉淀并如实呈现', 11.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P25 问卷设计 =================
s = new_slide()
header(s, '用户需求调研：问卷设计（节选）', 25)
qs = [
    ('Q1', '入学初期办理校园事务，平均需要多长时间找到准确指引？', '验证信息获取效率痛点'),
    ('Q2', '你通常通过哪些渠道获取办事信息？（多选）', '确定产品替代对象与渠道'),
    ('Q3', '选课时学长学姐的建议对你的决策影响有多大？（1—5 分）', '验证学业规划模块需求强度'),
    ('Q4', '当感到焦虑或压力大时，你更愿意向谁倾诉？', '验证心理疏导模块接受度'),
    ('Q5', '如果有一个 7×24 在线的"数字学长"，你最希望它解决哪三类问题？', '确定功能优先级'),
    ('Q6', '（面向学工教师）日常答疑类咨询占您工作时间比例约多少？', '验证 B 端减负价值'),
]
for i, (num, q, purpose) in enumerate(qs):
    y = 1.55 + i * 0.85
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, y, 0.75, 0.55, fill=BLUE, radius=0.3)
    add_text(s, 0.9, y, 0.75, 0.55, [([(num, 13, True, WHITE)], PP_ALIGN.CENTER, 0, 1.0)])
    add_text(s, 1.9, y + 0.02, 7.6, 0.6, [([(q, 12.5, False, DARK)], PP_ALIGN.LEFT, 0, 1.1)])
    add_text(s, 9.8, y + 0.02, 2.9, 0.6, [([(purpose, 11, False, GRAY)], PP_ALIGN.LEFT, 0, 1.1)])
add_text(s, 0.9, 6.75, 11.5, 0.5,
         [([('完整 8 题问卷与学工访谈提纲见商业计划书附录一 · 目标样本 ≥300 份 · 2026.09—10 执行', 11.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P26 财务预测 =================
s = new_slide()
header(s, '财务预测：三年收入与利润（万元）', 26)
add_shape(s, MSO_SHAPE.RECTANGLE, 9.0, 1.45, 0.22, 0.16, fill=BLUE)
add_text(s, 9.28, 1.38, 1.0, 0.3, [([('营业收入', 11, False, DARK)], PP_ALIGN.LEFT, 0, 1.0)])
add_shape(s, MSO_SHAPE.RECTANGLE, 10.15, 1.45, 0.22, 0.16, fill=MGRAY)
add_text(s, 10.43, 1.38, 1.0, 0.3, [([('成本费用', 11, False, DARK)], PP_ALIGN.LEFT, 0, 1.0)])
add_shape(s, MSO_SHAPE.OVAL, 11.32, 1.43, 0.2, 0.2, fill=ORANGE)
add_text(s, 11.58, 1.38, 1.0, 0.3, [([('净利润', 11, False, ORANGE)], PP_ALIGN.LEFT, 0, 1.0)])
years = ['第 1 年（2027）', '第 2 年（2028）', '第 3 年（2029）']
rev = [23.5, 86.3, 209.4]
cost = [26.0, 83.3, 190.1]
net = [-2.5, 3.0, 19.3]
BASE_Y, TOP_Y, MAXV = 6.05, 2.1, 220.0
SC = (BASE_Y - TOP_Y) / MAXV
gx = [2.0, 6.2, 10.4]
bar_w = 1.05
for i in range(3):
    add_text(s, gx[i] - 0.4, 6.2, 2.6, 0.35, [([(years[i], 12.5, True, DARK)], PP_ALIGN.CENTER, 0, 1.0)])
    h1 = rev[i] * SC
    h2 = cost[i] * SC
    add_shape(s, MSO_SHAPE.RECTANGLE, gx[i], BASE_Y - h1, bar_w, h1, fill=BLUE)
    add_shape(s, MSO_SHAPE.RECTANGLE, gx[i] + bar_w + 0.18, BASE_Y - h2, bar_w, h2, fill=MGRAY)
    add_text(s, gx[i] - 0.05, BASE_Y - h1 - 0.38, 1.15, 0.32,
             [([(f'{rev[i]:g}', 13, True, BLUE)], PP_ALIGN.CENTER, 0, 1.0)])
    add_text(s, gx[i] + 1.18, BASE_Y - h2 - 0.38, 1.15, 0.32,
             [([(f'{cost[i]:g}', 11, False, GRAY)], PP_ALIGN.CENTER, 0, 1.0)])
ny = [BASE_Y - v * SC for v in net]
for i in range(3):
    add_shape(s, MSO_SHAPE.OVAL, gx[i] + 1.05, ny[i] - 0.1, 0.2, 0.2, fill=ORANGE)
    add_text(s, gx[i] + 0.85, ny[i] - 0.38, 0.8, 0.32,
             [([(f'{net[i]:+.1f}', 12, True, ORANGE)], PP_ALIGN.CENTER, 0, 1.0)])
for i in range(2):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Inches(gx[i] + 1.15), Inches(ny[i]),
                                Inches(gx[i + 1] + 1.05), Inches(ny[i + 1]))
    ln.line.color.rgb = ORANGE
    ln.line.width = Pt(2.25)
    ln.shadow.inherit = False
add_text(s, 0.9, 1.45, 7.9, 0.5,
         [([('累计签约：4 所 → 16 所 → 40 所', 13, True, DARK)], PP_ALIGN.LEFT, 0, 1.0)])
add_text(s, 0.9, 6.65, 11.5, 0.5,
         [([('第 2 年盈亏平衡 · 第 3 年净利率 9.2%（成长期 SaaS 合理区间）· 第 4 年起向 15%—25% 演进', 12, False, GRAY)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P27 敏感性分析 =================
s = new_slide()
header(s, '敏感性分析：最坏情况仍可控（第 3 年）', 27)
rows, cols = 5, 4
tbl_shape = s.shapes.add_table(rows, cols, Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.4))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.3)
for j in range(1, 4):
    tbl.columns[j].width = Inches(2.67)
data = [
    ['情景', '关键变量变动', '第 3 年收入（万元）', '第 3 年净利润（万元）'],
    ['基准情景', '—', '209.4', '19.3'],
    ['保守情景 A', '签约学校数下降 30%', '约 168', '约 -8 至 +6（收缩弹性成本后）'],
    ['保守情景 B', '客单价下降 20%（结构恶化）', '约 167.5', '约 -23（触发止损预案）'],
    ['保守情景 C', '续费率由 80% 降至 60%', '约 194', '约 +3.5'],
]
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
            run_tf(run, data[r][c], 12.5, True, WHITE)
        elif r == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHTBLUE
            run_tf(run, data[r][c], 12, True if c == 0 else False, DARK)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            run_tf(run, data[r][c], 11.5, True if c == 0 else False, DARK if c == 0 else GRAY)
add_text(s, 0.9, 6.45, 11.5, 0.6,
         [([('结论：单一不利变动下项目仍接近盈亏平衡；弹性成本与签约节奏联动 + 6 个月现金储备兜底', 12.5, True, BLUE)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P28 融资需求 =================
s = new_slide()
header(s, '融资需求：种子轮 50 万元', 28)
add_text(s, 1.0, 2.0, 4.6, 2.5,
         [([('50', 60, True, BLUE), (' 万元', 24, True, BLUE)], PP_ALIGN.LEFT, 0, 1.0),
          ([('种子轮 · 拟出让股权 10%—15%', 15, True, DARK)], PP_ALIGN.LEFT, 10, 1.1),
          ([('启动时点：2027 年 6 月', 13, False, GRAY)], PP_ALIGN.LEFT, 4, 1.2),
          ([('（大创结题 + 试点成效数据确认后）', 13, False, GRAY)], PP_ALIGN.LEFT, 0, 1.2)])
add_text(s, 1.0, 5.0, 4.6, 1.4,
         [([('投资逻辑：', 13.5, True, BLUE)], PP_ALIGN.LEFT, 4, 1.2),
          ([('结题背书 + 试点数据支撑估值', 12.5, False, DARK)], PP_ALIGN.LEFT, 2, 1.2),
          ([('按 3—5 倍市销率，第 3 年估值约 630—1050 万元', 12.5, False, DARK)], PP_ALIGN.LEFT, 0, 1.2)])
chart_data = CategoryChartData()
chart_data.categories = ['产品研发', '试点部署', '市场拓展', '运营储备']
chart_data.add_series('资金用途', (20, 10, 12, 8))
gf = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(6.1), Inches(1.6), Inches(4.6), Inches(4.9), chart_data)
ch = gf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.RIGHT
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11)
ch.legend.font.color.rgb = DARK
plot = ch.plots[0]
plot.has_data_labels = True
dl = plot.data_labels
dl.show_percentage = True
dl.show_value = False
dl.number_format = '0%'
dl.number_format_is_linked = False
dl.font.size = Pt(11)
dl.font.color.rgb = WHITE
dl.font.bold = True
for pt, c in zip(ch.series[0].points, [BLUE, MIDBLUE, ORANGE, RGBColor(0xA6, 0xC5, 0xE0)]):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = c
add_text(s, 10.7, 2.2, 1.9, 3.6,
         [([('20 万（40%）', 12, True, BLUE)], PP_ALIGN.LEFT, 8, 1.1),
          ([('10 万（20%）', 12, True, MIDBLUE)], PP_ALIGN.LEFT, 8, 1.1),
          ([('12 万（24%）', 12, True, ORANGE)], PP_ALIGN.LEFT, 8, 1.1),
          ([('8 万（16%）', 12, True, RGBColor(0xA6, 0xC5, 0xE0))], PP_ALIGN.LEFT, 0, 1.1)])

# ================= P29 项目延伸 =================
s = new_slide()
header(s, '项目延伸：从产品到校园知识生态', 29)
exts = [
    ('产品延伸', 'V2.0 SaaS 多租户版本；增值服务（考研规划包、竞赛辅导资料库）', BLUE),
    ('生态延伸', '开放学长经验贡献入口，众包沉淀 + 审核机制 → 校园知识生态平台', ORANGE),
    ('模式延伸', '区域代理授权探索；从"卖软件"走向"运营校园知识生态"', MIDBLUE),
]
for i, (t, d, c) in enumerate(exts):
    x = 0.9 + i * 3.95
    cd = card(s, x, 1.7, 3.6, 3.3, fill=WHITE, line=c, line_w=1.5)
    chip(s, x + 1.28, 2.0, 1.04, 0.44, t[:2], c, size=13)
    add_text(s, x + 0.25, 2.7, 3.1, 2.1,
             [([(t, 15, True, DARK)], PP_ALIGN.CENTER, 4, 1.1),
              ([(d, 11.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.3)])
add_text(s, 0.9, 5.4, 11.5, 1.1,
         [([('三步走：', 13, True, BLUE), ('第 1—2 年订阅收入为主 → 第 2—3 年叠加增值服务 → 第 3 年起平台化与区域代理', 12.5, False, DARK)], PP_ALIGN.CENTER, 0, 1.2)])

# ================= P30 负责人+个人成长 =================
s = new_slide()
header(s, '团队（1/4）：负责人徐麟阁与个人成长', 30)
c1 = card(s, 0.9, 1.6, 5.7, 4.6, fill=LIGHTBLUE, line=None)
add_text(s, 1.25, 1.85, 5.0, 0.6, [([('徐麟阁 · 项目负责人 / 算法工程师', 16, True, BLUE)], PP_ALIGN.LEFT, 0, 1.0)])
add_text(s, 1.25, 2.6, 5.0, 3.4,
         [([('职责：', 12.5, True, DARK), ('总体统筹与进度管理；商业模式设计与融资对接', 12.5, False, DARK)], PP_ALIGN.LEFT, 3, 1.2),
          ([('产出：', 12.5, True, DARK), ('四层系统架构设计；RAG 问答链路搭建；财务模型与商业计划书', 12.5, False, DARK)], PP_ALIGN.LEFT, 3, 1.2),
          ([('个人成长：', 12.5, True, ORANGE)], PP_ALIGN.LEFT, 3, 1.2),
          ([('知识：', 12, True, DARK), ('课堂 NLP 理论 → RAG 工程落地（提示工程、检索优化）', 12, False, DARK)], PP_ALIGN.LEFT, 2, 1.2),
          ([('能力：', 12, True, DARK), ('完成实训作业 → 独立设计架构 → 牵头产品化与商业测算', 12, False, DARK)], PP_ALIGN.LEFT, 2, 1.2),
          ([('价值观：', 12, True, DARK), ('用 AI 传承朋辈经验，让"传帮带"数字化', 12, False, DARK)], PP_ALIGN.LEFT, 0, 1.2)])
c2 = card(s, 6.85, 1.6, 5.6, 4.6, fill=LIGHTORANGE, line=None)
add_text(s, 7.2, 1.85, 5.0, 0.6, [([('困难—帮扶—解决（真实经历）', 15, True, ORANGE)], PP_ALIGN.LEFT, 0, 1.0)])
add_text(s, 7.2, 2.6, 5.0, 3.4,
         [([('【技术困难案例：', 12, True, DARK), ('调试 RAG 检索时召回率长期不达标 → 在指导教师陈圣波建议下系统对比多种切分与嵌入策略 → 最终召回质量达标，掌握了检索链路优化的完整方法】', 12, False, DARK)], PP_ALIGN.LEFT, 6, 1.25),
          ([('【实践困难案例：', 12, True, DARK), ('调研问卷初版回收率低 → 依托创新创业学院渠道与辅导员协助改进投放方式 → 完成目标样本回收】', 12, False, DARK)], PP_ALIGN.LEFT, 6, 1.25),
          ([('提示：', 12, True, ORANGE), ('两个案例须替换为真实经历后再答辩，细节要经得起追问', 12, False, GRAY)], PP_ALIGN.LEFT, 0, 1.25)])

# ================= P31 团队成员+团队成长 =================
s = new_slide()
header(s, '团队（2/4）：团队成员与团队成长', 31)
add_text(s, 0.9, 1.45, 11.5, 0.5, [([('董芮铭 · 知识工程师 / 前端工程师', 15, True, BLUE)], PP_ALIGN.LEFT, 0, 1.0)])
add_text(s, 0.9, 1.95, 11.5, 1.0,
         [([('职责：知识图谱构建（Neo4j）与数据治理；跨平台客户端开发与交互设计', 12, False, DARK)], PP_ALIGN.LEFT, 0, 1.2),
          ([('成长：实体抽取练习 → 图谱 Schema 设计 → 独立完成客户端工程化', 12, False, GRAY)], PP_ALIGN.LEFT, 0, 1.2)])
add_text(s, 0.9, 3.15, 11.5, 0.5, [([('章宇洲 · 运营与市场', 15, True, ORANGE)], PP_ALIGN.LEFT, 0, 1.0)])
add_text(s, 0.9, 3.65, 11.5, 1.0,
         [([('职责：需求调研执行与内容运营；试点推广与客户成功；文档与合规管理', 12, False, DARK)], PP_ALIGN.LEFT, 0, 1.2),
          ([('成长：问卷设计 → 访谈执行 → 运营体系搭建', 12, False, GRAY)], PP_ALIGN.LEFT, 0, 1.2)])
ban = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 4.95, 11.5, 1.35, fill=LIGHTBLUE, radius=0.12)
shape_text(ban, [([('团队成长：', 13.5, True, BLUE), ('三人小队磨合 → 一人多角分工 → 模块牵头人制，协作效率与项目质量同步提升', 12.5, False, DARK)], PP_ALIGN.CENTER, 3, 1.2),
                 ([('贡献台账（2026 新增评审项）：每人对应模块产出均可逐项核查，与研发时间线匹配', 11.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.2)])

# ================= P32 专家顾问 =================
s = new_slide()
header(s, '团队（3/4）：专家顾问（行业 / 商业 / 学术技术）', 32)
tag(s, 11.9, 1.15, '待确认')
advs = [
    ('行业顾问', '【姓名 / 单位】\n教育信息化行业经验\n对接高校客户与渠道资源', BLUE),
    ('商业顾问', '【姓名 / 单位】\nSaaS 商业模式与融资\n辅导财务规划与股权设计', ORANGE),
    ('学术技术顾问', '【姓名 / 单位】\nNLP / 知识图谱方向\n算法方案评审与前沿把关', MIDBLUE),
]
for i, (t, d, c) in enumerate(advs):
    x = 0.9 + i * 3.95
    cd = card(s, x, 1.7, 3.6, 3.3, fill=WHITE, line=c, line_w=1.5)
    chip(s, x + 1.28, 2.0, 1.04, 0.44, t[:2], c, size=13)
    add_text(s, x + 0.25, 2.7, 3.1, 2.1,
             [([(t, 15, True, DARK)], PP_ALIGN.CENTER, 4, 1.1)] +
             [([(ln, 11.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.3) for ln in d.split('\n')])
add_text(s, 0.9, 5.4, 11.5, 1.2,
         [([('说明：顾问团队均为行业、商业、学术技术方向（按评审要求不放本校教师）；', 12, False, GRAY)], PP_ALIGN.CENTER, 2, 1.25),
          ([('指导教师陈圣波与学校资源见下一页"学校给予项目的支撑"', 12, False, GRAY)], PP_ALIGN.CENTER, 0, 1.25)])

# ================= P33 学校支撑 =================
s = new_slide()
header(s, '团队（4/4）：学校给予项目的支撑', 33)
sups = [
    ('指导教师', '陈圣波——技术路线把关、算法方案评审、高校资源对接', BLUE),
    ('创新创业学院', '创业辅导、场地支持、政策对接；渠道引荐（2026.12 前合作备忘录目标）', ORANGE),
    ('课程与实训基础', '机器学习、NLP 课程 + 三项实训项目，直接转化为核心技术能力', MIDBLUE),
]
for i, (t, d, c) in enumerate(sups):
    x = 0.9 + i * 3.95
    cd = card(s, x, 1.7, 3.6, 3.3, fill=WHITE, line=c, line_w=1.5)
    chip(s, x + 1.28, 2.0, 1.04, 0.44, t[:2], c, size=13)
    add_text(s, x + 0.25, 2.7, 3.1, 2.1,
             [([(t, 15, True, DARK)], PP_ALIGN.CENTER, 4, 1.1),
              ([(d, 11.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.3)])
add_text(s, 0.9, 5.4, 11.5, 1.2,
         [([('学创融合：项目创意产生、技术研发、落地实践与课程培养环节一一对应，体现"以赛促学"育人成效', 13, True, BLUE)], PP_ALIGN.CENTER, 0, 1.2)])

# ================= P34 媒体报道 =================
s = new_slide()
header(s, '媒体报道与传播计划（计划）', 34)
tag(s, 11.9, 1.15, '计划')
meds = [
    ('校内传播', '官方公众号"学长说"系列专栏；嵌入新生入学教育场景', BLUE),
    ('行业露出', '高教会 / 教育信息化展会案例输出（2027）', ORANGE),
    ('媒体报道目标', '2027 年内 1—2 篇（区域教育媒体 / 行业自媒体）', MIDBLUE),
    ('内容资产', '试点成效白皮书：问答量、减负时长、满意度数据', BLUE),
]
for i, (t, d, c) in enumerate(meds):
    x = 0.9 + (i % 2) * 5.95
    y = 1.7 + (i // 2) * 2.3
    card(s, x, y, 5.6, 1.9, fill=WHITE, line=c, line_w=1.5)
    add_shape(s, MSO_SHAPE.RECTANGLE, x, y, 0.12, 1.9, fill=c)
    add_text(s, x + 0.4, y + 0.2, 5.0, 1.5,
             [([(t, 16, True, c)], PP_ALIGN.LEFT, 4, 1.1),
              ([(d, 12, False, DARK)], PP_ALIGN.LEFT, 0, 1.25)])
add_text(s, 0.9, 6.45, 11.5, 0.5,
         [([('说明：以上为传播计划；实际媒体报道形成后将如实补充报道链接与版面', 11.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P35 社会价值 =================
s = new_slide()
header(s, '社会价值：就业 · 行业 · 家国', 35)
vals = [
    ('就业带动', '公司化运营提供全职与实习岗位；校园实习生计划吸纳低年级学生', BLUE),
    ('行业进步', '为高校智能导学提供可复制方案；推动朋辈教育数字化成为行业方向', ORANGE),
    ('家国情怀', '服务立德树人：用技术把校园经验代际传承，让教育公平延伸到信息服务', MIDBLUE),
]
for i, (t, d, c) in enumerate(vals):
    x = 0.9 + i * 3.95
    cd = card(s, x, 1.7, 3.6, 3.3, fill=WHITE, line=c, line_w=1.5)
    chip(s, x + 1.28, 2.0, 1.04, 0.44, t[:2], c, size=13)
    add_text(s, x + 0.25, 2.7, 3.1, 2.1,
             [([(t, 15, True, DARK)], PP_ALIGN.CENTER, 4, 1.1),
              ([(d, 11.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.3)])
add_text(s, 0.9, 5.4, 11.5, 1.2,
         [([('项目使命：', 13, True, BLUE), ('让每一段大学时光，都有智慧与温暖相伴——用 AI 补位朋辈互助，让育人资源普惠每一位学生', 12.5, False, DARK)], PP_ALIGN.CENTER, 0, 1.2)])

# ================= P36 发展规划 =================
s = new_slide()
header(s, '发展规划：三年三步走', 36)
line_y = 5.35
add_shape(s, MSO_SHAPE.RECTANGLE, 1.4, line_y - 0.02, 10.6, 0.045, fill=MGRAY)
phases = [
    ('验证期', '2026.09 — 2027.06', '校内跑通闭环\nMVP · 成效报告', '签约 1—4 所', BLUE),
    ('放量期', '2027.07 — 2028.12', '公司化运营\n渠道合作 · 盈亏平衡', '累计 16 所 · 年收入约 86 万', ORANGE),
    ('平台期', '2029.01 — 2029.12', '规模化复制\nSaaS 多租户 · 增值服务', '累计 40 所 · 年收入约 210 万', MIDBLUE),
]
for i, (name, t, body, mile, c) in enumerate(phases):
    x = 1.0 + i * 3.95
    card(s, x, 1.6, 3.6, 2.5, fill=WHITE, line=c, line_w=1.5)
    chip(s, x + 1.05, 1.85, 1.5, 0.45, name, c, size=13)
    add_text(s, x + 0.3, 2.5, 3.0, 0.45, [([(t, 13, True, DARK)], PP_ALIGN.CENTER, 0, 1.0)])
    add_text(s, x + 0.3, 3.0, 3.0, 0.9, [([(ln, 11.5, False, GRAY)], PP_ALIGN.CENTER, 0, 1.2) for ln in body.split('\n')])
    add_shape(s, MSO_SHAPE.OVAL, x + 1.55, line_y - 0.16, 0.34, 0.34, fill=c, line=WHITE, line_w=1.5)
    add_text(s, x, 5.55, 3.6, 0.7, [([(mile, 14, True, c)], PP_ALIGN.CENTER, 0, 1.1)])
add_text(s, 0.9, 6.5, 11.5, 0.6,
         [([('衔接大创周期：2026.09 立项 → 2027.06 结题（校内全量部署 + 首份成效报告）→ 公司化运营', 12.5, False, DARK)], PP_ALIGN.CENTER, 0, 1.0)])

# ================= P37 风险与应对 =================
s = new_slide()
header(s, '风险与应对：Top 3', 37)
risks = [
    ('市场风险', '采购周期长、预算不确定', '学院级轻量入口 · 标杆案例背书 · 6 个月现金安全线', BLUE, LIGHTBLUE),
    ('技术风险', '大模型“幻觉”损害口碑', 'RAG 约束生成 + 来源引用 · 白名单答案库 · 知识审核机制', ORANGE, LIGHTORANGE),
    ('数据合规', '学生隐私与安全审查', '私有化部署 · 最小化采集 · 脱敏训练 · 配合校方审查', MIDBLUE, LIGHTBLUE),
]
for i, (t, r, m, c, fill) in enumerate(risks):
    x = 0.9 + i * 3.95
    card(s, x, 1.8, 3.6, 3.5, fill=fill, line=None)
    add_text(s, x + 0.25, 2.1, 3.1, 0.6, [([(t, 18, True, c)], PP_ALIGN.CENTER, 0, 1.0)])
    add_text(s, x + 0.25, 2.8, 3.1, 0.8, [([(r, 12.5, False, DARK)], PP_ALIGN.CENTER, 0, 1.2)])
    add_shape(s, MSO_SHAPE.RECTANGLE, x + 1.3, 3.7, 1.0, 0.03, fill=c)
    add_text(s, x + 0.25, 3.9, 3.1, 1.2, [([('应对：', 12.5, True, c), (m, 12, False, DARK)], PP_ALIGN.CENTER, 0, 1.25)])
add_text(s, 0.9, 5.6, 11.5, 1.15,
         [([('其余风险（政策/经营）应对详见商业计划书表 10-1；', 12, False, GRAY),
            ('整体原则：弹性成本与签约进度联动，保守扩张', 12, True, DARK)], PP_ALIGN.CENTER, 2, 1.1),
          ([('2026 规则红线：材料真实性一票否决——本项目数据保守测算并如实标注假设，知识产权由团队持有', 11, True, ORANGE)], PP_ALIGN.CENTER, 0, 1.1)])

# ================= P38 结语 =================
s = new_slide()
add_shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.334, 7.5, fill=BLUE)
add_shape(s, MSO_SHAPE.RECTANGLE, 0, 7.38, 13.334, 0.12, fill=ORANGE)
add_text(s, 0.9, 2.35, 11.5, 1.2,
         [([('让每一段大学时光，都有智慧与温暖相伴', 34, True, WHITE)], PP_ALIGN.CENTER, 0, 1.0)])
add_shape(s, MSO_SHAPE.RECTANGLE, 5.77, 3.95, 1.8, 0.05, fill=ORANGE)
add_text(s, 0.9, 4.25, 11.5, 0.6,
         [([('从"数字学长"出发，构建可传承的校园知识生态', 17, False, RGBColor(0xD9, 0xE2, 0xF3))], PP_ALIGN.CENTER, 0, 1.0)])
add_text(s, 0.9, 5.6, 11.5, 1.0,
         [([('恳请各位老师批评指正 · 谢谢聆听', 20, True, RGBColor(0xF4, 0xB1, 0x83))], PP_ALIGN.CENTER, 0, 1.0)])

prs.save(OUT)
print('已生成:', OUT, f'（共 {TOTAL} 页）')
